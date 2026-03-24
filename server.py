from flask import Flask, request, send_file
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor as PPTXColor
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Pt as DocxPt
from docx.shared import RGBColor as DocxColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
import fitz  # PyMuPDF
import io
from deep_translator import GoogleTranslator

app = Flask(__name__)
CORS(app)

def translate_text(text):
    if not text or len(text.strip()) < 3 or text.isnumeric(): return None
    try:
        return GoogleTranslator(source='auto', target='ar').translate(text)
    except: return None

# 1. محرك البوربوينت (تنسيق متناسق للملفات الكبيرة)
@app.route('/translate-pptx', methods=['POST'])
def handle_pptx():
    file = request.files['file']
    print(f"\n🚀 جاري معالجة البوربوينت بأحجام متناسقة: {file.filename}")
    prs = Presentation(file)
    
    for slide in prs.slides:
        slide_text = "".join([shape.text for shape in slide.shapes if hasattr(shape, "text_frame")])
        char_count = len(slide_text.strip())
        
        # أحجام "متناسقة" عشان تناسب كل الملفات
        if char_count < 250:      
            en_size, ar_size = Pt(22), Pt(18) # بدلاً من 26/20
        elif char_count < 600:    
            en_size, ar_size = Pt(16), Pt(13) # بدلاً من 18/14
        else:                     
            en_size, ar_size = Pt(12), Pt(10) # للملفات المزدحمة جداً
            
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text.strip():
                for paragraph in shape.text_frame.paragraphs:
                    original = paragraph.text.strip()
                    if len(original) < 3: continue
                    translated = translate_text(original)
                    if translated:
                        paragraph.text = ""
                        paragraph.alignment = PP_ALIGN.LEFT
                        
                        # الإنجليزي
                        run_en = paragraph.add_run()
                        run_en.text = original + "   " # 3 فراغات للتوازن
                        run_en.font.size = en_size
                        run_en.font.bold = True
                        run_en.font.color.rgb = PPTXColor(60, 60, 60)
                        
                        # العربي
                        run_ar = paragraph.add_run()
                        run_ar.text = translated
                        run_ar.font.size = ar_size
                        run_ar.font.bold = True
                        run_ar.font.color.rgb = PPTXColor(29, 185, 84)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return send_file(output, as_attachment=True, download_name=f"KR_Pro_{file.filename}")

# 2. محرك الوورد (يحافظ على التنسيق الأصلي)
@app.route('/translate-docx', methods=['POST'])
def handle_docx():
    file = request.files['file']
    print(f"\n📝 جاري معالجة الوورد: {file.filename}")
    doc = Document(file)

    def process_p(p):
        original = p.text.strip()
        if len(original) > 3 and not original.isnumeric():
            translated = translate_text(original)
            if translated:
                p.text = ""
                run_en = p.add_run(original + "   ")
                run_en.font.size = DocxPt(11)
                run_en.font.color.rgb = DocxColor(60, 60, 60)
                
                run_ar = p.add_run(translated)
                run_ar.font.size = DocxPt(11)
                run_ar.font.bold = True
                run_ar.font.color.rgb = DocxColor(29, 185, 84)

    for p in doc.paragraphs: process_p(p)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs: process_p(p)

    output = io.BytesIO()
    doc.save(output)
    output.seek(0)
    return send_file(output, as_attachment=True, download_name=f"KR_Pro_{file.filename}")

# 3. محرك الـ PDF (جداول وورد)
@app.route('/translate-pdf', methods=['POST'])
def handle_pdf():
    file = request.files['file']
    doc_pdf = fitz.open(stream=file.read(), filetype="pdf")
    doc_word = Document()
    for page in doc_pdf:
        blocks = page.get_text("blocks")
        blocks.sort(key=lambda b: (b[1], b[0]))
        table = doc_word.add_table(rows=0, cols=2)
        table.style = 'Table Grid'
        for b in blocks:
            original = b[4].strip()
            if len(original) > 5:
                translated = translate_text(original)
                if translated:
                    row = table.add_row().cells
                    row[0].text = original
                    p_ar = row[1].paragraphs[0]
                    p_ar.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                    run_ar = p_ar.add_run(translated)
                    run_ar.font.color.rgb = DocxColor(29, 185, 84)
                    run_ar.font.bold = True
    output = io.BytesIO()
    doc_word.save(output)
    output.seek(0)
    return send_file(output, as_attachment=True, download_name=f"KR_Pro_{file.filename.replace('.pdf', '.docx')}")

if __name__ == '__main__':
    import os
    port = int(os.environ.get("PORT", 5000))
    app.run(host='0.0.0.0', port=port)